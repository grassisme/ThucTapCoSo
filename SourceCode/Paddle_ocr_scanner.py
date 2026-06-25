import os
import sys
import json
import logging
import argparse
from pathlib import Path

os.environ["PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK"] = "True"


def check_dependencies():
    missing = []
    try:
        import paddleocr
    except ImportError:
        missing.append("paddleocr")
    try:
        from PIL import Image
    except ImportError:
        missing.append("Pillow")
    if missing:
        print(f"[!] Thiếu thư viện: {', '.join(missing)}")
        print(f"    Cài đặt: pip install {' '.join(missing)}")
        sys.exit(1)


class PaddleOCRScanner:
    def __init__(self, lang: str = "vi", use_gpu: bool = False):
        from paddleocr import PaddleOCR
        logging.getLogger("ppocr").setLevel(logging.ERROR)
        logging.getLogger("paddle").setLevel(logging.ERROR)

        init_ok = False

        if not init_ok:
            try:
                self.ocr = PaddleOCR(lang=lang, device="gpu" if use_gpu else "cpu")
                init_ok = True
                self._version = "v3_device"
            except Exception:
                pass

        if not init_ok:
            try:
                self.ocr = PaddleOCR(lang=lang)
                init_ok = True
                self._version = "v3_minimal"
            except Exception:
                pass

        if not init_ok:
            self.ocr = PaddleOCR(
                use_angle_cls=True,
                lang=lang,
                use_gpu=use_gpu,
                show_log=False,
            )
            self._version = "v2"

    def _run_ocr(self, image_path: str):
        try:
            return self.ocr.ocr(image_path, cls=True)
        except TypeError:
            pass
        try:
            return self.ocr.ocr(image_path)
        except Exception:
            pass
        return list(self.ocr.predict(image_path))

    def _parse_result(self, result) -> list:
        if not result:
            return []
        if not isinstance(result, list):
            result = list(result)

        lines = []
        for item in result:
            if isinstance(item, dict):
            # ✅ Format PaddleOCR v3: dict có rec_texts + rec_scores
                texts = item.get("rec_texts", [])
                scores = item.get("rec_scores", [])
                for t, s in zip(texts, scores):
                    if t and t.strip():
                        lines.append((t.strip(), s))

            elif isinstance(item, list):
            # Format v2: [[box], (text, conf)]
                for line in item:
                    try:
                        val = line[1]
                        if isinstance(val, (tuple, list)):
                            lines.append((val[0], val[1]))
                        elif isinstance(val, str):
                            lines.append((val, 1.0))
                    except Exception:
                        pass
        return lines

    def ocr_image(self, image_path: str) -> str:
        result = self._run_ocr(image_path)
        parsed = self._parse_result(result)
        return "\n".join(text for text, _ in parsed if text)

    def ocr_image_detail(self, image_path: str) -> list:
        result = self._run_ocr(image_path)
        parsed = self._parse_result(result)
        return [{"text": t, "confidence": round(c, 4)} for t, c in parsed]


# ─── EASYOCR (đọc được dấu tiếng Việt) ───────────────────────────────────────
class EasyOCRScanner:
    """Engine EasyOCR — hỗ trợ tiếng Việt đầy đủ. Cài: pip install easyocr"""
    _reader_cache = None   # (langs_tuple, reader)

    def __init__(self, lang: str = "vi"):
        self._langs = ["vi", "en"] if str(lang).lower().startswith("vi") else ["en"]

    def _get_reader(self):
        import easyocr
        key = tuple(self._langs)
        if (EasyOCRScanner._reader_cache is None
                or EasyOCRScanner._reader_cache[0] != key):
            print(f"[*] Đang load EasyOCR {self._langs} ...")
            EasyOCRScanner._reader_cache = (key, easyocr.Reader(self._langs, gpu=False))
            print("[✓] EasyOCR sẵn sàng")
        return EasyOCRScanner._reader_cache[1]

    def ocr_image(self, image_path: str) -> str:
        reader = self._get_reader()
        lines = reader.readtext(image_path, detail=0, paragraph=True)
        return "\n".join(l for l in lines if l and str(l).strip())

    def ocr_image_detail(self, image_path: str) -> list:
        reader = self._get_reader()
        rows = reader.readtext(image_path, detail=1, paragraph=False)
        return [{"text": t, "confidence": round(float(c), 4)}
                for _box, t, c in rows if t and str(t).strip()]


# ─── TESSERACT (cần cài binary hệ thống + gói ngôn ngữ vie) ───────────────────
class TesseractScanner:
    """Engine Tesseract. Cài: pip install pytesseract  +  Tesseract-OCR (gói 'vie').
    Nếu binary không nằm trong PATH, đặt biến môi trường TESSERACT_CMD trỏ tới
    tesseract.exe (vd: C:\\Program Files\\Tesseract-OCR\\tesseract.exe)."""

    def __init__(self, lang: str = "vie"):
        self._lang = "vie+eng" if str(lang).lower().startswith("vi") else "eng"

    @staticmethod
    def _configure():
        import pytesseract
        cmd = os.environ.get("TESSERACT_CMD")
        if cmd:
            pytesseract.pytesseract.tesseract_cmd = cmd
        return pytesseract

    @staticmethod
    def available() -> bool:
        try:
            pt = TesseractScanner._configure()
            pt.get_tesseract_version()
            return True
        except Exception:
            return False

    def ocr_image(self, image_path: str) -> str:
        pt = self._configure()
        from PIL import Image
        return pt.image_to_string(Image.open(image_path), lang=self._lang).strip()

    def ocr_image_detail(self, image_path: str) -> list:
        txt = self.ocr_image(image_path)
        return [{"text": l, "confidence": 1.0} for l in txt.splitlines() if l.strip()]


# ─── FACTORY + SO SÁNH OCR ───────────────────────────────────────────────────
def make_scanner(engine: str = "auto", lang: str = "vi"):
    """Tạo scanner theo engine. 'auto': tiếng Việt -> EasyOCR, còn lại -> PaddleOCR."""
    engine = (engine or "auto").lower()
    if engine == "auto":
        engine = "easyocr" if str(lang).lower().startswith("vi") else "paddle"
    if engine == "easyocr":
        return EasyOCRScanner(lang)
    if engine == "tesseract":
        return TesseractScanner(lang)
    return PaddleOCRScanner(lang=lang)


def _vi_diacritic_ratio(text: str) -> float:
    """% ký tự mang dấu riêng của tiếng Việt / tổng số chữ cái — proxy chất lượng OCR VN."""
    vi = ("ăâđêôơưĂÂĐÊÔƠƯáàảãạấầẩẫậắằẳẵặéèẻẽẹếềểễệíìỉĩị"
          "óòỏõọốồổỗộớờởỡợúùủũụứừửữựýỳỷỹỵ")
    letters = [c for c in (text or "") if c.isalpha()]
    if not letters:
        return 0.0
    return round(sum(1 for c in text if c in vi) / len(letters) * 100, 1)


def compare_ocr(image_path: str, lang: str = "vi") -> list:
    """Chạy nhiều engine OCR trên cùng một ảnh để so sánh (phục vụ đánh giá mục 1.3).
    Trả về list dict: {engine, name, text, chars, vi_ratio, seconds, available, error}."""
    import time
    engines = [("paddle", "PaddleOCR"), ("easyocr", "EasyOCR"), ("tesseract", "Tesseract")]
    out = []
    for key, name in engines:
        t0 = time.perf_counter()
        text, err, avail = "", "", True
        try:
            text = make_scanner(key, lang).ocr_image(image_path) or ""
        except ImportError as e:
            err, avail = f"chưa cài ({e})", False
        except Exception as e:
            err, avail = str(e), False
        out.append({
            "engine": key, "name": name, "text": text, "chars": len(text),
            "vi_ratio": _vi_diacritic_ratio(text),
            "seconds": round(time.perf_counter() - t0, 2),
            "available": avail, "error": err,
        })
    return out


def scan_image(path: str, scanner) -> str:
    return scanner.ocr_image(path)


def scan_pdf(path: str, scanner: PaddleOCRScanner) -> str:
    try:
        import fitz
    except ImportError:
        print("[!] Cần cài PyMuPDF: pip install PyMuPDF")
        return ""

    import tempfile
    doc = fitz.open(path)
    all_text = []

    with tempfile.TemporaryDirectory() as tmp:
        for page_num, page in enumerate(doc):
            mat = fitz.Matrix(300 / 72, 300 / 72)
            pix = page.get_pixmap(matrix=mat)
            img_path = os.path.join(tmp, f"page_{page_num + 1}.png")
            pix.save(img_path)
            text = scanner.ocr_image(img_path)
            if text.strip():
                all_text.append(f"=== Trang {page_num + 1} ===\n{text}")
            print(f"  [PDF] Trang {page_num + 1}/{len(doc)}")

    doc.close()
    return "\n\n".join(all_text)


def scan_docx(path: str, scanner: PaddleOCRScanner) -> str:
    try:
        from docx import Document
    except ImportError:
        print("[!] Cần cài python-docx: pip install python-docx")
        return ""

    import tempfile, io
    from PIL import Image

    doc = Document(path)
    parts = []

    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    for table in doc.tables:
        for row in table.rows:
            row_texts = [c.text.strip() for c in row.cells if c.text.strip()]
            if row_texts:
                parts.append(" | ".join(row_texts))

    with tempfile.TemporaryDirectory() as tmp:
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                img_path = os.path.join(tmp, f"embed_{rel.rId}.png")
                try:
                    img = Image.open(io.BytesIO(rel.target_part.blob)).convert("RGB")
                    img.save(img_path)
                    t = scanner.ocr_image(img_path)
                    if t.strip():
                        parts.append(f"[Ảnh nhúng]\n{t}")
                except Exception:
                    pass

    return "\n".join(parts)


def scan_xlsx(path: str) -> str:
    try:
        import openpyxl
    except ImportError:
        print("[!] Cần cài openpyxl: pip install openpyxl")
        return ""

    wb = openpyxl.load_workbook(path, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"=== Sheet: {sheet_name} ===")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def scan_pptx(path: str, scanner: PaddleOCRScanner) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        print("[!] Cần cài python-pptx: pip install python-pptx")
        return ""

    import tempfile, io
    from PIL import Image

    prs = Presentation(path)
    parts = []

    for i, slide in enumerate(prs.slides):
        slide_texts = [f"=== Slide {i + 1} ==="]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
            if shape.shape_type == 13:
                with tempfile.TemporaryDirectory() as tmp:
                    img_path = os.path.join(tmp, f"s{i}_{shape.shape_id}.png")
                    try:
                        img = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                        img.save(img_path)
                        t = scanner.ocr_image(img_path)
                        if t.strip():
                            slide_texts.append(f"[Ảnh]\n{t}")
                    except Exception:
                        pass
        parts.append("\n".join(slide_texts))

    return "\n\n".join(parts)


def scan_txt(path: str) -> str:
    for enc in ("utf-8", "utf-8-sig", "utf-16", "cp1258", "latin-1"):
        try:
            with open(path, encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, LookupError):
            continue
    return ""


SUPPORTED_EXTENSIONS = {
    ".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp",
    ".pdf", ".docx", ".xlsx", ".xls", ".pptx", ".txt",
}


def scan_file(path: str, scanner: PaddleOCRScanner) -> str:
    ext = Path(path).suffix.lower()
    if ext in {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp"}:
        return scan_image(path, scanner)
    elif ext == ".pdf":
        return scan_pdf(path, scanner)
    elif ext == ".docx":
        return scan_docx(path, scanner)
    elif ext in {".xlsx", ".xls"}:
        return scan_xlsx(path)
    elif ext == ".pptx":
        return scan_pptx(path, scanner)
    elif ext == ".txt":
        return scan_txt(path)
    else:
        print(f"[!] Không hỗ trợ: {ext} | Hỗ trợ: {', '.join(sorted(SUPPORTED_EXTENSIONS))}")
        return ""


def quick_scan(file_path: str, lang: str = "vi", engine: str = "auto") -> str:
    """
    Dùng trong code:
        from Paddle_ocr_scanner import quick_scan
        text = quick_scan(r"G:\\path\\to\\file.jpg", lang="vi", engine="auto")

    engine: "auto" (vi -> EasyOCR, khác -> PaddleOCR) | "paddle" | "easyocr" | "tesseract".
    Nếu engine chọn chưa cài, tự lùi về PaddleOCR và in hướng dẫn cài đặt.
    """
    check_dependencies()
    eng = (engine or "auto").lower()
    if eng == "auto":
        eng = "easyocr" if str(lang).lower().startswith("vi") else "paddle"
    try:
        return scan_file(file_path, make_scanner(eng, lang))
    except ImportError as e:
        print(f"[!] Engine '{eng}' chưa sẵn sàng ({e}).")
        if eng == "easyocr":
            print("    Cài: pip install easyocr")
        elif eng == "tesseract":
            print("    Cài: pip install pytesseract  +  Tesseract-OCR (gói ngôn ngữ 'vie')")
        print("[!] Tạm dùng PaddleOCR (có thể rụng dấu tiếng Việt).")
        return scan_file(file_path, PaddleOCRScanner(lang=lang))


def main():
    parser = argparse.ArgumentParser(description="Scan file bằng PaddleOCR")
    parser.add_argument("files", nargs="+", help="Đường dẫn file")
    parser.add_argument("--lang", default="vi")
    parser.add_argument("--gpu", action="store_true")
    parser.add_argument("--output", "-o")
    parser.add_argument("--json", action="store_true")
    parser.add_argument("--detail", action="store_true")
    args = parser.parse_args()

    check_dependencies()
    scanner = PaddleOCRScanner(lang=args.lang, use_gpu=args.gpu)
    print(f"[*] PaddleOCR {scanner._version}")

    results = {}
    for file_path in args.files:
        if not os.path.exists(file_path):
            print(f"[!] Không tìm thấy: {file_path}")
            continue
        print(f"\n[>] {file_path}")
        if args.detail and Path(file_path).suffix.lower() in {".png",".jpg",".jpeg",".bmp",".tiff",".tif",".webp"}:
            detail = scanner.ocr_image_detail(file_path)
            results[file_path] = detail
            for item in detail:
                print(f"  [{item['confidence']:.2f}] {item['text']}")
        else:
            text = scan_file(file_path, scanner)
            results[file_path] = text
            print(text or "  (Không có nội dung)")

    if args.output:
        if args.json or args.output.endswith(".json"):
            with open(args.output, "w", encoding="utf-8") as f:
                json.dump(results, f, ensure_ascii=False, indent=2)
        else:
            with open(args.output, "w", encoding="utf-8") as f:
                for p, c in results.items():
                    f.write(f"{'='*60}\nFile: {p}\n{'='*60}\n")
                    if isinstance(c, list):
                        for item in c:
                            f.write(f"[{item['confidence']:.2f}] {item['text']}\n")
                    else:
                        f.write(c + "\n\n")
        print(f"[✓] Đã lưu: {args.output}")


if __name__ == "__main__":
    main()